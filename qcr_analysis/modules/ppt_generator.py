# -*- coding: utf-8 -*-
"""
=============================================================================
PPT生成模块
=============================================================================
负责生成PowerPoint报告
=============================================================================
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from typing import Any, Dict, Optional

import sys
sys.path.append(str(Path(__file__).parent.parent))
from config import PPT_STYLE
from modules.llm_service import LLMService, LLMGenerationError


class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self, template_path: Optional[Path] = None):
        """
        初始化PPT生成器
        
        Args:
            template_path: PPT模板文件路径（可选）
        """
        if template_path and template_path.exists():
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.font_name = PPT_STYLE['font_name']
        self.title_font_size = PPT_STYLE['title_font_size']
        self.subtitle_font_size = PPT_STYLE['subtitle_font_size']
        self.body_font_size = PPT_STYLE['body_font_size']
    
    def add_textbox(self, slide, left, top, width, height, text,
                   font_name=None, font_size=None, bold=False):
        """
        在幻灯片中添加文本框
        
        Args:
            slide: 幻灯片对象
            left, top, width, height: 位置和尺寸（Inches对象）
            text: 文本内容
            font_name: 字体名称
            font_size: 字号
            bold: 是否加粗
            
        Returns:
            文本框对象
        """
        if font_name is None:
            font_name = self.font_name
        if font_size is None:
            font_size = self.body_font_size
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        return textbox
    
    def add_image(self, slide, image_path: str, left, top, width=None, height=None):
        """
        在幻灯片中添加图片
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            left, top: 位置（Inches对象）
            width, height: 尺寸（Inches对象，可选）
        """
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    
    def build_homepage(self, payload: Dict[str, Any]):
        """
        生成首页幻灯片
        
        Args:
            payload: 数据载荷
        """
        slide_layout = self.prs.slide_layouts[5]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # ========== 标题区域 ==========
        title_text = "落入D等级 数据汇总分析结果"
        self.add_textbox(
            slide,
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8),
            title_text,
            font_size=self.title_font_size,
            bold=True
        )
        
        # ========== 正文区域 ==========
        bullet_left = Inches(0.5)
        bullet_top = Inches(1.2)
        bullet_box = slide.shapes.add_textbox(bullet_left, bullet_top, Inches(9), Inches(2.5))
        frame = bullet_box.text_frame
        frame.word_wrap = True
        frame.clear()
        
        # 提取数据
        start_str, end_str = payload.get("coverage_period", ("-", "-"))
        week_start, week_end = payload.get("week_range", ("-", "-"))
        product_list = payload.get("unique_models", [])
        product_text = "、".join(product_list[:4]) if product_list else "暂无"
        total_records = payload.get("total_records", 0)
        
        # Bullet 1: 基本信息
        p1 = frame.add_paragraph()
        p1.text = (
            f"{week_start}-{week_end}共收到落入D等级产品数据{total_records}条，"
            f"覆盖周期为{start_str}-{end_str}，"
            f"产品为{product_text}等，共计{len(product_list)}款。"
        )
        p1.level = 0
        p1.font.name = self.font_name
        p1.font.size = Pt(self.body_font_size)
        
        # Bullet 2: 审核原因占比
        reason_df = payload.get("reason_stats", pd.DataFrame())
        if not reason_df.empty:
            row_map = {row["审核原因"]: f"{row['占比']}%" for _, row in reason_df.iterrows()}
            p2 = frame.add_paragraph()
            p2.text = (
                f"审核原因中，7天无理由占比{row_map.get('7天无理由', '0%')}，"
                f"15天质量换新占比{row_map.get('15天质量换新', '0%')}，"
                f"质量维修占比{row_map.get('质量维修', '0%')}，"
                f"180天只换不修占比{row_map.get('180天只换不修', '0%')}。"
            )
            p2.level = 0
            p2.font.name = self.font_name
            p2.font.size = Pt(self.body_font_size)
        
        # Bullet 3: 七天无理由机型
        model_7d_df = payload.get("model_7d_dist", pd.DataFrame())
        if not model_7d_df.empty:
            top_items = model_7d_df.head(4)
            parts = [f"{row['机型名称']}占比{row['占比']}%" for _, row in top_items.iterrows()]
            p3 = frame.add_paragraph()
            p3.text = "七天无理由中，" + "，".join(parts) + "。"
            p3.level = 0
            p3.font.name = self.font_name
            p3.font.size = Pt(self.body_font_size)
        
        # Bullet 4: 非七天无理由机型
        model_non7d_df = payload.get("model_non_7d_dist", pd.DataFrame())
        if not model_non7d_df.empty:
            top_items = model_non7d_df.head(4)
            parts = [f"{row['机型名称']}占比{row['占比']}%" for _, row in top_items.iterrows()]
            p4 = frame.add_paragraph()
            p4.text = "非七天无理由中，" + "，".join(parts) + "。"
            p4.level = 0
            p4.font.name = self.font_name
            p4.font.size = Pt(self.body_font_size)
        
        # ========== 图表区域 ==========
        # 3张饼图横向排列
        chart_y = Inches(4.2)
        chart_width = Inches(2.5)
        
        self.add_image(slide, payload.get("reason_chart_path"), Inches(1.0), chart_y, width=chart_width)
        self.add_image(slide, payload.get("model_7d_chart_path"), Inches(4.2), chart_y, width=chart_width)
        self.add_image(slide, payload.get("model_non7d_chart_path"), Inches(7.4), chart_y, width=chart_width)
        
        print("✓ PPT首页已生成")
    
    def build_detail_slide(self, model_name: str, suffix: str, entry: Dict[str, Any],
                          use_llm: bool = False, llm_service: Optional[LLMService] = None,
                          llm_params: Optional[Dict] = None):
        """
        生成详情页幻灯片
        
        Args:
            model_name: 机型名称
            suffix: 分类后缀
            entry: 数据条目
            use_llm: 是否使用LLM
            llm_service: LLM服务对象
            llm_params: LLM参数
        """
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # ========== 标题区域 ==========
        title = f"{model_name} {suffix}分类"
        self.add_textbox(
            slide,
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8),
            title,
            font_size=self.subtitle_font_size,
            bold=True
        )
        
        # ========== 正文区域 ==========
        print(f"→ 生成详情页：机型='{model_name}', 类型='{suffix}'")
        text_content = ""
        
        if use_llm and llm_service:
            try:
                print(f"   调用Kimi生成观点中...")
                llm_params = llm_params or {}
                text_content = llm_service.generate_summary(
                    entry.get("category_df", pd.DataFrame()),
                    timeout=llm_params.get("timeout", 60),
                    top_n=llm_params.get("top_n", 3),
                    coverage_threshold=llm_params.get("coverage", 80.0),
                    focus_threshold=llm_params.get("focus", 10.0)
                )
                print(f"   ✓ Kimi返回观点")
            except LLMGenerationError as exc:
                print(f"   ✗ LLM生成失败: {exc}")
                text_content = LLMService.get_fallback_text(
                    entry.get("clean_model", ""),
                    suffix,
                    entry.get("total_records", 0)
                )
                print(f"   → 使用本地模板观点")
        else:
            text_content = LLMService.get_fallback_text(
                entry.get("clean_model", ""),
                suffix,
                entry.get("total_records", 0)
            )
            print(f"   （未启用LLM）使用本地模板观点")
        
        # 添加正文文本框
        self.add_textbox(
            slide,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(2.8),
            text_content,
            font_size=self.body_font_size
        )
        
        # ========== 图表区域 ==========
        # 柱状图居中显示
        chart_width = Inches(4.5)
        chart_left = Inches(2.75)
        chart_top = Inches(4.5)
        
        self.add_image(slide, entry.get("chart_path"), chart_left, chart_top, width=chart_width)
        print(f"← 完成：机型='{model_name}', 类型='{suffix}' 的详情页\n")
    
    def build_top_issue_summary_slide(self, top_issue_result: Dict[str, Any]):
        """
        生成Top Issue总结页
        
        Args:
            top_issue_result: Top Issue分析结果
        """
        slide_layout = self.prs.slide_layouts[5]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        top_n = top_issue_result.get('top_n', 10)
        total_records = top_issue_result.get('total_records', 0)
        issue_stats = top_issue_result.get('issue_stats', pd.DataFrame())
        
        # 标题
        title_text = f"Top {top_n} Issue 分布总结"
        self.add_textbox(
            slide,
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8),
            title_text,
            font_size=self.subtitle_font_size,
            bold=True
        )
        
        # 数据概览
        if not issue_stats.empty:
            cumulative_pct = issue_stats['累计占比(%)'].iloc[-1]
            overview_text = (
                f"数据总量：{total_records} 条记录\n"
                f"Top {len(issue_stats)} Issue累计占比：{cumulative_pct:.2f}%"
            )
        else:
            overview_text = f"数据总量：{total_records} 条记录"
        
        self.add_textbox(
            slide,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(0.8),
            overview_text,
            font_size=self.body_font_size
        )
        
        # 添加Top Issue总览图
        summary_chart_path = top_issue_result.get('summary_chart_path')
        if summary_chart_path and Path(summary_chart_path).exists():
            self.add_image(
                slide,
                summary_chart_path,
                Inches(1.0), Inches(2.2),
                width=Inches(8.0)
            )
        
        print(f"✓ 生成Top Issue总结页")
    
    def build_top_issue_detail_slide(self, issue_detail: Dict[str, Any],
                                     use_llm: bool = False,
                                     llm_service: Optional[LLMService] = None,
                                     llm_params: Optional[Dict] = None):
        """
        生成Top Issue详情页
        
        Args:
            issue_detail: Issue详情数据
            use_llm: 是否使用LLM
            llm_service: LLM服务对象
            llm_params: LLM参数
        """
        slide_layout = self.prs.slide_layouts[5]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        issue_name = issue_detail.get('issue_name', '')
        count = issue_detail.get('count', 0)
        percentage = issue_detail.get('percentage', 0)
        rank = issue_detail.get('rank', 0)
        model_dist = issue_detail.get('model_dist', pd.DataFrame())
        chart_path = issue_detail.get('chart_path')
        
        # 标题
        title_text = f"Issue #{rank} - {issue_name} ({count}台, {percentage:.1f}%)"
        self.add_textbox(
            slide,
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.6),
            title_text,
            font_size=self.subtitle_font_size,
            bold=True
        )
        
        # 左侧：机型分布图
        left_section_left = Inches(0.5)
        left_section_top = Inches(1.0)
        left_section_width = Inches(4.5)
        
        if chart_path and Path(chart_path).exists():
            self.add_image(
                slide,
                chart_path,
                left_section_left, left_section_top,
                width=left_section_width
            )
        
        # 右侧：LLM质量管理分析
        right_section_left = Inches(5.2)
        right_section_top = Inches(1.0)
        right_section_width = Inches(4.3)
        right_section_height = Inches(6.0)
        
        analysis_text = ""
        
        if use_llm and llm_service:
            try:
                print(f"   调用LLM分析Issue: {issue_name}...")
                analysis_text = llm_service.analyze_top_issue(
                    issue_name=issue_name,
                    count=count,
                    percentage=percentage,
                    model_dist=model_dist,
                    timeout=llm_params.get('timeout', 60) if llm_params else 60
                )
                print(f"   ✓ LLM分析完成")
            except Exception as e:
                print(f"   ✗ LLM分析失败: {e}")
                analysis_text = self._get_fallback_issue_analysis(
                    issue_name, count, percentage, model_dist
                )
                print(f"   → 使用本地模板分析")
        else:
            analysis_text = self._get_fallback_issue_analysis(
                issue_name, count, percentage, model_dist
            )
            print(f"   （未启用LLM）使用本地模板分析")
        
        self.add_textbox(
            slide,
            right_section_left, right_section_top,
            right_section_width, right_section_height,
            analysis_text,
            font_size=12
        )
        
        print(f"← 完成Issue详情页：{issue_name}\n")
    
    def _get_fallback_issue_analysis(self, issue_name: str, count: int,
                                     percentage: float, model_dist: pd.DataFrame) -> str:
        """
        生成默认的Issue分析文本（当LLM不可用时）
        
        Args:
            issue_name: Issue名称
            count: 数量
            percentage: 占比
            model_dist: 机型分布DataFrame
            
        Returns:
            分析文本
        """
        lines = []
        lines.append(f"【问题特征】")
        lines.append(f"{issue_name}是本批次的重要质量问题。")
        lines.append("")
        
        lines.append(f"【影响范围】")
        lines.append(f"共影响{count}台设备，占比{percentage:.2f}%。")
        lines.append("")
        
        lines.append(f"【机型关联分析】")
        if not model_dist.empty:
            top_model = model_dist.iloc[0]
            model_count = len(model_dist)
            lines.append(f"涉及{model_count}款机型。")
            lines.append(f"主要集中在{top_model['机型名称']} ")
            lines.append(f"({top_model['数量']}台, {top_model['占比(%)']:.1f}%)。")
        lines.append("")
        
        lines.append(f"【可能原因】")
        lines.append(f"1. 产品设计或工艺问题")
        lines.append(f"2. 供应链质量波动")
        lines.append(f"3. 使用环境因素")
        lines.append("")
        
        lines.append(f"【改进建议】")
        lines.append(f"1. 加强质量检测和管控")
        lines.append(f"2. 优化供应商管理")
        lines.append(f"3. 改进产品设计和工艺")
        
        return "\n".join(lines)
    
    def build_top_issue_slides(self, top_issue_result: Dict[str, Any],
                              use_llm: bool = False,
                              llm_service: Optional[LLMService] = None,
                              llm_params: Optional[Dict] = None):
        """
        生成所有Top Issue相关页面
        
        Args:
            top_issue_result: Top Issue分析结果
            use_llm: 是否使用LLM
            llm_service: LLM服务对象
            llm_params: LLM参数
        """
        # 1. 生成总结页
        self.build_top_issue_summary_slide(top_issue_result)
        
        # 2. 生成每个Issue的详情页
        issue_details = top_issue_result.get('issue_details', [])
        for issue_detail in issue_details:
            self.build_top_issue_detail_slide(
                issue_detail,
                use_llm, llm_service, llm_params
            )
        
        print(f"✓ 完成{len(issue_details)}个Issue详情页")
    
    def generate(self, summary_payload: Dict[str, Any], output_path: Path,
                use_llm: bool = False, llm_service: Optional[LLMService] = None,
                llm_params: Optional[Dict] = None):
        """
        生成完整的PPT报告
        
        Args:
            summary_payload: 汇总数据载荷
            output_path: 输出文件路径
            use_llm: 是否使用LLM
            llm_service: LLM服务对象
            llm_params: LLM参数
        """
        print("\n" + "="*60)
        print("开始生成PPT报告")
        print("="*60)
        
        # 生成首页
        print("\n生成首页...")
        self.build_homepage(summary_payload)
        
        # 生成详情页
        model_details = summary_payload.get("model_details", {})
        for model_name, entries in model_details.items():
            for suffix in ["7天无理由", "非7天无理由"]:
                entry = entries.get(suffix)
                if not entry:
                    continue
                self.build_detail_slide(
                    model_name, suffix, entry,
                    use_llm, llm_service, llm_params
                )
        
        # 生成Top Issue分析页面（如果有）
        top_issue_result = summary_payload.get("top_issue_result")
        if top_issue_result:
            print("\n生成Top Issue分析页面...")
            self.build_top_issue_slides(
                top_issue_result,
                use_llm, llm_service, llm_params
            )
        
        # 保存PPT
        self.prs.save(str(output_path))
        print("\n" + "="*60)
        print(f"✅ PPT报告已生成：{output_path}")
        print("="*60 + "\n")

