# -*- coding: utf-8 -*-
"""报告生成服务"""
import pandas as pd
from pathlib import Path
from typing import Dict, Optional
from datetime import datetime

import sys
sys.path.append(str(Path(__file__).parent.parent))
from modules.ppt_generator import PPTGenerator
from modules.llm_service import LLMService
from prompts import (
    TOP_ISSUE_SUMMARY_PROMPT,
    TOP_MODEL_OVERVIEW_PROMPT,
)

# PPT样式辅助函数
def set_title_style(shape, text=""):
    """设置标题样式：微软雅黑，28号，居中"""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    
    if hasattr(shape, 'text_frame'):
        tf = shape.text_frame
    else:
        tf = shape
    
    if text:
        tf.text = text
    
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = '微软雅黑'
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True

def set_body_style(shape, text="", font_size=11):
    """设置正文样式：微软雅黑，指定号数，左对齐"""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    
    if hasattr(shape, 'text_frame'):
        tf = shape.text_frame
    else:
        tf = shape
    
    if text:
        tf.text = text
    
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = '微软雅黑'
        paragraph.font.size = Pt(font_size)

def set_table_style(table):
    """设置表格样式：微软雅黑，10号"""
    from pptx.util import Pt
    
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = '微软雅黑'
                paragraph.font.size = Pt(10)

class ReportService:
    """报告生成服务"""
    
    def __init__(self, output_dir: str or Path):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_weekly_ppt(self, payload, batch_name, template_path=None, use_llm=False, llm_config=None):
        """生成Weekly Report PPT"""
        ppt_filename = f"weekly_report_{batch_name}.pptx"
        ppt_path = self.output_dir / ppt_filename
        
        ppt_generator = PPTGenerator(template_path=template_path, use_llm=use_llm, llm_config=llm_config)
        ppt_generator.generate_complete_ppt(summary_payload=payload, output_path=str(ppt_path))
        
        return ppt_path
    
    def generate_top_issue_ppt(self, payload, batch_name, template_path=None, use_llm=False, llm_config=None):
        """生成Top Issue PPT（重构版：标题28号，正文11号，表格10号，概览页整合）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        ppt_filename = f"top_issue_report_{batch_name}.pptx"
        ppt_path = self.output_dir / ppt_filename

        prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        issue_stats = payload.get("issue_stats")
        issue_details = payload.get("issue_details", [])
        summary_chart = payload.get("summary_chart")
        total_records = payload.get("total_records")
        top_n = payload.get("top_n")
        # 预计算 AI 总览
        ai_overview = None
        if use_llm and issue_stats is not None and len(issue_stats) > 0:
            try:
                ai_overview = self._summarize_issue_overview(issue_stats, llm_config)
            except Exception:
                ai_overview = None

        # 【页1：情况概览 - 整合所有内容】
        overview_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        
        # 标题（28号，居中）
        title_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        set_title_style(title_tb, f"Top Issue 分析报告 - {batch_name}")
        
        # 统计信息（11号，左对齐）
        stats_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(3), Inches(0.8))
        set_body_style(stats_tb, f"分析数据：{total_records:,} 条\nIssue分类：{top_n} 个", 11)
        
        # Top Issue表格（10号，左侧）
        if issue_stats is not None and len(issue_stats) > 0:
            table_rows = min(len(issue_stats), 10) + 1
            table_cols = 4
            try:
                table_shape = overview_slide.shapes.add_table(
                    table_rows, table_cols, 
                    Inches(0.5), Inches(2.0), 
                    Inches(4.2), Inches(3.0)
                )
                table = table_shape.table
                headers = ["排名", "Issue", "数量", "占比%"]
                for c, h in enumerate(headers):
                    table.cell(0, c).text = str(h)
                for r in range(1, table_rows):
                    row_data = issue_stats.iloc[r-1]
                    table.cell(r, 0).text = str(row_data["排名"])
                    table.cell(r, 1).text = str(row_data["Issue名称"])[:20]
                    table.cell(r, 2).text = str(row_data["数量"])
                    table.cell(r, 3).text = str(row_data["占比(%)"])
                set_table_style(table)
            except Exception:
                pass
        
        # 总览图（右侧，自适应）
        if summary_chart and Path(summary_chart).exists():
            try:
                overview_slide.shapes.add_picture(str(summary_chart), Inches(5.0), Inches(1.5), width=Inches(4.5))
            except Exception:
                pass
        
        # AI总结（11号，左对齐，底部）
        if ai_overview:
            ai_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
            set_body_style(ai_tb, f"📊 AI洞察：{ai_overview}", 11)

        
        # 【页2-11：Issue详情页】（前10个Issue）
        for detail in issue_details[:10]:
            detail_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
            
            # 标题（28号，居中）
            title_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            set_title_style(title_tb, f"Issue #{detail['rank']}: {detail['issue_name']}")
            
            # 统计信息（11号，左对齐）
            stats_text = f"数量：{detail['count']} ({detail['percentage']}%)\n机型数：{detail['model_count']}"
            stats_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.6))
            set_body_style(stats_tb, stats_text, 11)
            
            # 机型分布图（自适应大小）
            chart_path = detail.get("chart_path")
            if chart_path and Path(chart_path).exists():
                try:
                    detail_slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.8), width=Inches(9))
                except Exception:
                    pass
            
            # AI洞察（11号，左对齐）
            if use_llm:
                try:
                    model_dist = detail.get("model_distribution")
                    if model_dist is not None and len(model_dist) > 0:
                        ai_insight = self._summarize_issue_detail(detail['issue_name'], model_dist, llm_config)
                        insight_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
                        set_body_style(insight_tb, f"💡 AI洞察：{ai_insight}", 11)
                except Exception:
                    pass
        
        prs.save(str(ppt_path))
        return ppt_path
    
    def _summarize_model_llm(self, model_name, category_dist, llm_config):
        """对单个机型的分类分布使用LLM生成摘要（简化版）"""
        try:
            llm = LLMService(
                api_key=llm_config.get("api_key") if llm_config else None,
                api_url=llm_config.get("api_url") if llm_config else None,
                model=llm_config.get("model") if llm_config else None,
            )
            top_rows = category_dist.head(10)
            lines = ["分类\t数量\t占比"]
            for _, row in top_rows.iterrows():
                lines.append(f"{row['分类']}\t{row['数量']}\t{row['占比(%)']}")
            table_text = "\n".join(lines)
            
            prompt = (
                f"你是PQM质量专家，请分析机型'{model_name}'的问题分类分布，"
                f"用一句话总结主要问题集中在哪些分类，不超过50字。数据表：\n"
                f"{table_text}"
            )
            return llm.call_api([{"role": "user", "content": prompt}], timeout=int(llm_config.get("timeout", 60)) if llm_config else 60)
        except Exception as e:
            return f"AI 摘要生成失败: {e}"

    def _summarize_issue_overview(self, issue_stats, llm_config):
        """Top Issue 概览页一句话总结：分布+重点关注"""
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top_rows = issue_stats.head(5)
        lines = ["分类\t数量\t占比"]
        for _, row in top_rows.iterrows():
            lines.append(f"{row['Issue名称']}\t{row['数量']}\t{row['占比(%)']}")
        table_text = "\n".join(lines)
        prompt = (
            "你是PQM质量专家，请用一句中文总结当前Top Issue分布，点名重点关注的Issue和机型（如有信息），"
            "语气简洁，不超过50字。数据表：\n"
            f"{table_text}"
        )
        return llm.call_api([{"role": "user", "content": prompt}], timeout=int(llm_config.get("timeout", 60)) if llm_config else 60)

    def _summarize_issue_detail(self, issue_name, model_dist, llm_config):
        """针对单个Issue的机型分布，给出PQM视角洞察"""
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top_rows = model_dist.head(10)
        lines = ["机型\t数量\t占比"]
        for _, row in top_rows.iterrows():
            lines.append(f"{row['机型名称']}\t{row['数量']}\t{row['占比(%)']}")
        table_text = "\n".join(lines)
        prompt = (
            "你是PQM质量专家，请针对该Issue的机型分布给出一句洞察，不超过50字，"
            "可强调高风险机型或拦截建议。Issue: {issue}, 数据表：\n"
            f"{table_text}"
        ).format(issue=issue_name)
        return llm.call_api([{"role": "user", "content": prompt}], timeout=int(llm_config.get("timeout", 60)) if llm_config else 60)
    
    def _summarize_model_overview(self, top_models, llm_config):
        """Top Model 概览页一句话总结：问题复杂度分析"""
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top_rows = top_models.head(5)
        lines = ["机型\t分类数量\t记录数"]
        for _, row in top_rows.iterrows():
            lines.append(f"{row['机型名称']}\t{row['分类数']}\t{row['记录数']}")
        table_text = "\n".join(lines)
        prompt = (
            "你是PQM质量专家，请用一句中文总结Top Model的问题复杂度分布，指出哪些机型问题最复杂，"
            "语气简洁，不超过50字。数据表：\n"
            f"{table_text}"
        )
        return llm.call_api([{"role": "user", "content": prompt}], timeout=int(llm_config.get("timeout", 60)) if llm_config else 60)

    def generate_top_model_ppt(self, payload, batch_name, template_path=None, use_llm=False, llm_config=None):
        """生成Top Model PPT（重构版：标题28号，正文11号，表格10号，概览页整合）"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        ppt_filename = f"top_model_report_{batch_name}.pptx"
        ppt_path = self.output_dir / ppt_filename

        prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        model_stats = payload.get("model_stats")
        top_models = payload.get("top_models")
        model_details = payload.get("model_details", [])
        overall_chart = payload.get("overall_chart")
        comparison_chart = payload.get("comparison_chart")
        total_records = payload.get("total_records")
        total_models = payload.get("total_models")
        top_n = payload.get("top_n")

        # 预计算 AI 总览
        ai_overview = None
        if use_llm and top_models is not None and len(top_models) > 0:
            try:
                ai_overview = self._summarize_model_overview(top_models, llm_config)
            except Exception:
                ai_overview = None

        # 【页1：情况概览 - 整合所有内容】
        overview_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        
        # 标题（28号，居中）
        title_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
        set_title_style(title_tb, f"Top Model 分析报告 - {batch_name}\n(基于分类数量)")
        
        # 统计信息（11号，左对齐）
        stats_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(3.5), Inches(0.8))
        set_body_style(stats_tb, f"分析数据：{total_records:,} 条\n机型总数：{total_models} 个\nTop N：{top_n}", 11)
        
        # Top Model表格（10号，左侧）
        if top_models is not None and len(top_models) > 0:
            table_rows = min(len(top_models), 10) + 1
            table_cols = 5
            try:
                table_shape = overview_slide.shapes.add_table(
                    table_rows, table_cols, 
                    Inches(0.5), Inches(2.0), 
                    Inches(4.5), Inches(3.0)
                )
                table = table_shape.table
                headers = ["排名", "机型", "类别数", "记录数", "平均/类"]
                for c, h in enumerate(headers):
                    table.cell(0, c).text = str(h)
                for r in range(1, table_rows):
                    row_data = top_models.iloc[r-1]
                    table.cell(r, 0).text = str(row_data["排名"])
                    table.cell(r, 1).text = str(row_data["机型名称"])
                    table.cell(r, 2).text = str(row_data["分类数"])
                    table.cell(r, 3).text = str(row_data["记录数"])
                    table.cell(r, 4).text = str(row_data["平均每类记录数"])
                set_table_style(table)
            except Exception:
                pass
        
        # 整体分布图或对比图（右侧，自适应）
        chart_to_show = overall_chart or comparison_chart
        if chart_to_show and Path(chart_to_show).exists():
            try:
                overview_slide.shapes.add_picture(str(chart_to_show), Inches(5.2), Inches(1.2), width=Inches(4.3))
            except Exception:
                pass
        
        # AI总结（11号，左对齐，底部）
        if ai_overview:
            ai_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
            set_body_style(ai_tb, f"📊 AI洞察：{ai_overview}", 11)
        
        # 【页2-11：Model详情页】（前10个机型）
        for detail in model_details[:10]:
            detail_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
            
            # 标题（28号，居中）
            title_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            set_title_style(title_tb, f"机型 #{detail['rank']}: {detail['model_name']}")
            
            # 统计信息 + Top 5分类（11号，左对齐）
            category_dist = detail.get("category_distribution")
            top5_text = ""
            if category_dist is not None and len(category_dist) > 0:
                top5_text = "\n\nTop 5 分类：\n"
                for i, row in category_dist.head(5).iterrows():
                    top5_text += f"{i+1}. {row['分类']} - {row['数量']} ({row['占比(%)']}%)\n"
            
            stats_text = (
                f"分类数：{detail['category_count']}\n"
                f"记录数：{detail['total_records']}"
                f"{top5_text}"
            )
            stats_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(2.0))
            set_body_style(stats_tb, stats_text, 11)
            
            # 分类分布图（自适应大小）
            chart_path = detail.get("chart_path")
            if chart_path and Path(chart_path).exists():
                try:
                    detail_slide.shapes.add_picture(str(chart_path), Inches(0.5), Inches(3.2), width=Inches(9))
                except Exception:
                    pass
            
            # AI解读（11号，左对齐）
            if use_llm:
                try:
                    category_dist = detail.get("category_distribution")
                    if category_dist is not None and len(category_dist) > 0:
                        ai_insight = self._summarize_model_llm(detail['model_name'], category_dist, llm_config)
                        insight_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
                        set_body_style(insight_tb, f"💡 AI解读：{ai_insight}", 11)
                except Exception:
                    pass
        
        prs.save(str(ppt_path))
        return ppt_path

def create_report_service(output_dir):
    return ReportService(output_dir)

def generate_weekly_report(payload, output_dir, batch_name, template_path=None, use_llm=False, llm_config=None):
    service = ReportService(output_dir)
    return service.generate_weekly_ppt(payload, batch_name, template_path, use_llm, llm_config)

def generate_top_issue_report(payload, output_dir, batch_name, template_path=None, use_llm=False, llm_config=None):
    service = ReportService(output_dir)
    return service.generate_top_issue_ppt(payload, batch_name, template_path, use_llm, llm_config)

def generate_top_model_report(payload, output_dir, batch_name, template_path=None, use_llm=False, llm_config=None):
    service = ReportService(output_dir)
    return service.generate_top_model_ppt(payload, batch_name, template_path, use_llm, llm_config)

