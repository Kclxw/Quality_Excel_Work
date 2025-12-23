# -*- coding: utf-8 -*-
"""
重构后的PPT生成方法
将以下两个方法替换到report_service.py中的ReportService类
"""

def generate_top_issue_ppt_refactored(self, payload, batch_name, template_path=None, use_llm=False, llm_config=None):
    """生成Top Issue PPT（重构版：标题28号，正文11号，表格10号，概览页整合）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    ppt_filename = f"top_issue_report_{batch_name}.pptx"
    ppt_path = self.output_dir / ppt_filename

    prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    issue_stats = payload.get("issue_stats")
    issue_details = payload.get("issue_details", [])
    summary_chart = payload.get("summary_chart")
    total_records = payload.get("total_records")
    top_n = payload.get("top_n")

    # 预计算 AI 总览
    ai_overview = None
    if use_llm and issue_stats is not None and len(issue_stats) > 0:
        try:
            ai_overview = self._summarize_issue_overview(issue_stats, llm_config)
        except Exception:
            ai_overview = None

    # 【页1：概览页 - 整合所有内容】
    overview_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    
    # 标题
    title = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    set_title_style(title, f"Top Issue 分析报告 - {batch_name}")
    
    # 统计信息（左上）
    stats_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(3), Inches(0.8))
    set_body_style(stats_tb, f"总记录数：{total_records:,}\nIssue分类数：{top_n}", 11)
    
    # Top Issue表格（左侧）
    if issue_stats is not None and len(issue_stats) > 0:
        table_rows = min(len(issue_stats), 10) + 1
        table_cols = 4
        try:
            table_shape = overview_slide.shapes.add_table(
                table_rows, table_cols, 
                Inches(0.5), Inches(2.0), 
                Inches(4.2), Inches(3.0)
            )
            table = table_shape.table
            
            # 表头
            headers = ["排名", "Issue", "数量", "占比%"]
            for c, h in enumerate(headers):
                cell = table.cell(0, c)
                cell.text = str(h)
            
            # 数据行
            for r in range(1, table_rows):
                row_data = issue_stats.iloc[r-1]
                table.cell(r, 0).text = str(row_data["排名"])
                table.cell(r, 1).text = str(row_data["Issue名称"])[:20]  # 限制长度
                table.cell(r, 2).text = str(row_data["数量"])
                table.cell(r, 3).text = str(row_data["占比(%)"])
            
            # 设置表格样式
            set_table_style(table)
        except Exception as e:
            # 表格失败兜底
            tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.2), Inches(3.0))
            set_body_style(tb, f"Top Issue列表（表格失败：{e}）", 10)
    
    # 总览图（右侧）
    if summary_chart and Path(summary_chart).exists():
        try:
            overview_slide.shapes.add_picture(
                str(summary_chart), 
                Inches(5.0), Inches(1.5), 
                width=Inches(4.5)
            )
        except Exception:
            pass
    
    # AI总结（底部）
    if ai_overview:
        ai_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
        set_body_style(ai_tb, f"📊 AI洞察：{ai_overview}", 11)
    
    # 【页2-11：Issue详情页】（共10页）
    for detail in issue_details[:10]:
        detail_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        
        # 标题
        title_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        set_title_style(title_tb, f"Issue #{detail['rank']}: {detail['issue_name']}")
        
        # 统计信息
        stats_text = f"数量：{detail['count']} ({detail['percentage']}%)\n机型数：{detail['model_count']}"
        stats_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.6))
        set_body_style(stats_tb, stats_text, 11)
        
        # 机型分布图
        chart_path = detail.get("chart_path")
        if chart_path and Path(chart_path).exists():
            try:
                detail_slide.shapes.add_picture(
                    str(chart_path), 
                    Inches(0.5), Inches(1.8), 
                    width=Inches(9)
                )
            except Exception:
                pass
        
        # AI洞察（每个Issue）
        if use_llm:
            try:
                model_dist = detail.get("model_distribution")
                if model_dist is not None and len(model_dist) > 0:
                    ai_insight = self._summarize_issue_detail(
                        detail['issue_name'], 
                        model_dist, 
                        llm_config
                    )
                    insight_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
                    set_body_style(insight_tb, f"💡 AI洞察：{ai_insight}", 11)
            except Exception:
                pass
    
    prs.save(str(ppt_path))
    return ppt_path


def generate_top_model_ppt_refactored(self, payload, batch_name, template_path=None, use_llm=False, llm_config=None):
    """生成Top Model PPT（重构版：标题28号，正文11号，表格10号，概览页整合）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    ppt_filename = f"top_model_report_{batch_name}.pptx"
    ppt_path = self.output_dir / ppt_filename

    prs = Presentation(template_path) if template_path and Path(template_path).exists() else Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    model_stats = payload.get("model_stats")
    top_models = payload.get("top_models")
    model_details = payload.get("model_details", [])
    overall_chart = payload.get("overall_chart")
    comparison_chart = payload.get("comparison_chart")
    total_records = payload.get("total_records")
    total_models = payload.get("total_models")
    top_n = payload.get("top_n")

    # 预计算 AI 总览
    ai_overview = None
    if use_llm and top_models is not None and len(top_models) > 0:
        try:
            ai_overview = self._summarize_model_overview(top_models, llm_config)
        except Exception:
            ai_overview = None

    # 【页1：概览页 - 整合所有内容】
    overview_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
    
    # 标题
    title = overview_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    set_title_style(title, f"Top Model 分析报告 - {batch_name}\n基于问题类别数量")
    
    # 统计信息（左上）
    stats_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(3.5), Inches(0.8))
    set_body_style(stats_tb, f"总记录数：{total_records:,}\n机型总数：{total_models}\nTop N：{top_n}", 11)
    
    # Top Model表格（左侧）
    if top_models is not None and len(top_models) > 0:
        table_rows = min(len(top_models), 10) + 1
        table_cols = 5
        try:
            table_shape = overview_slide.shapes.add_table(
                table_rows, table_cols, 
                Inches(0.5), Inches(2.0), 
                Inches(4.5), Inches(3.0)
            )
            table = table_shape.table
            
            # 表头
            headers = ["排名", "机型", "类别数", "记录数", "平均/类"]
            for c, h in enumerate(headers):
                table.cell(0, c).text = str(h)
            
            # 数据行
            for r in range(1, table_rows):
                row_data = top_models.iloc[r-1]
                table.cell(r, 0).text = str(row_data["排名"])
                table.cell(r, 1).text = str(row_data["机型名称"])[:15]
                table.cell(r, 2).text = str(row_data["问题类别数"])
                table.cell(r, 3).text = str(row_data["记录数"])
                table.cell(r, 4).text = str(row_data["平均每类记录数"])
            
            # 设置表格样式
            set_table_style(table)
        except Exception as e:
            tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.5), Inches(3.0))
            set_body_style(tb, f"Top Model列表（表格失败：{e}）", 10)
    
    # 整体分布图或对比图（右侧）
    chart_to_show = overall_chart or comparison_chart
    if chart_to_show and Path(chart_to_show).exists():
        try:
            overview_slide.shapes.add_picture(
                str(chart_to_show), 
                Inches(5.2), Inches(1.2), 
                width=Inches(4.3)
            )
        except Exception:
            pass
    
    # AI总结（底部）
    if ai_overview:
        ai_tb = overview_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
        set_body_style(ai_tb, f"📊 AI洞察：{ai_overview}", 11)
    
    # 【页2-11：Model详情页】（共10页）
    for detail in model_details[:10]:
        detail_slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        
        # 标题
        title_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        set_title_style(title_tb, f"机型 #{detail['rank']}: {detail['model_name']}")
        
        # 统计信息
        stats_text = (
            f"问题类别数：{detail['category_count']}\n"
            f"记录数：{detail['total_records']}\n"
            f"7天无理由：{detail['return_7day_count']} ({detail['return_7day_pct']}%)\n"
            f"质量问题：{detail['quality_count']} ({detail['quality_pct']}%)"
        )
        stats_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4), Inches(1.2))
        set_body_style(stats_tb, stats_text, 11)
        
        # 分类分布图
        chart_path = detail.get("chart_path")
        if chart_path and Path(chart_path).exists():
            try:
                detail_slide.shapes.add_picture(
                    str(chart_path), 
                    Inches(0.5), Inches(2.4), 
                    width=Inches(9)
                )
            except Exception:
                pass
        
        # AI解读（每个机型）
        if use_llm:
            try:
                category_dist = detail.get("category_distribution")
                if category_dist is not None and len(category_dist) > 0:
                    ai_insight = self._summarize_model_llm(
                        detail['model_name'], 
                        category_dist, 
                        llm_config
                    )
                    insight_tb = detail_slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
                    set_body_style(insight_tb, f"💡 AI解读：{ai_insight}", 11)
            except Exception:
                pass
    
    prs.save(str(ppt_path))
    return ppt_path


# 辅助方法（需要添加到ReportService类）
def _summarize_issue_overview(self, issue_stats, llm_config):
    """总览级别AI摘要（一句话）"""
    try:
        from modules.llm_service import LLMService
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top3 = issue_stats.head(3)
        summary = ", ".join([f"{row['Issue名称']}({row['占比(%)']}%)" for _, row in top3.iterrows()])
        prompt = f"请用一句话总结Top Issue分布：{summary}。重点关注哪些Issue和机型？"
        return llm.call_api([{"role": "user", "content": prompt}], timeout=30)
    except Exception as e:
        return f"AI总结失败: {e}"

def _summarize_issue_detail(self, issue_name, model_dist, llm_config):
    """单Issue详情AI洞察（从PQM视角）"""
    try:
        from modules.llm_service import LLMService
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top5 = model_dist.head(5)
        models = ", ".join([f"{row['机型名称']}({row['占比(%)']}%)" for _, row in top5.iterrows()])
        prompt = f"从PQM视角，用一句话总结Issue '{issue_name}' 的机型分布：{models}。质量管理建议？"
        return llm.call_api([{"role": "user", "content": prompt}], timeout=30)
    except Exception as e:
        return f"洞察失败: {e}"

def _summarize_model_overview(self, top_models, llm_config):
    """总览级别AI摘要（一句话）"""
    try:
        from modules.llm_service import LLMService
        llm = LLMService(
            api_key=llm_config.get("api_key") if llm_config else None,
            api_url=llm_config.get("api_url") if llm_config else None,
            model=llm_config.get("model") if llm_config else None,
        )
        top3 = top_models.head(3)
        summary = ", ".join([f"{row['机型名称']}({row['问题类别数']}类)" for _, row in top3.iterrows()])
        prompt = f"请用一句话总结Top Model分布：{summary}。问题复杂度分析？"
        return llm.call_api([{"role": "user", "content": prompt}], timeout=30)
    except Exception as e:
        return f"AI总结失败: {e}"

