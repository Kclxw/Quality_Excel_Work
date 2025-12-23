# -*- coding: utf-8 -*-
"""
=============================================================================
QCR 数据分析与 PPT 自动生成工具
=============================================================================

功能概述:
---------
1. 读取Excel数据并进行数据库去重
2. 指定处理时间周期（基于日期列）
3. 读取MTM表格并映射机型名称
4. 统计四种审核原因（7天无理由、15天质量换新、180天只换不修、质量维修）
5. 统计7天无理由/非7天无理由的机型分布
6. 按机型统计分类描述词频次
7. 为每个机型的所有分类生成详细数据文件
8. 生成表格+饼图/柱状图
9. 可选：生成PPT报告（首页+详情页）
10. 可选：调用Kimi LLM生成智能分析摘要

依赖安装:
---------
pip install pandas openpyxl matplotlib pymysql sqlalchemy python-pptx requests

配置说明:
---------
【重要】Kimi API密钥配置（使用--use-llm时必需）：
  方式1（推荐）：直接在代码第171行修改 DEFAULT_KIMI_API_KEY 常量
  方式2：设置环境变量 KIMI_API_KEY（优先级高于代码配置）

可选环境变量（如不设置则使用默认值）：
- KIMI_API_URL: Kimi API地址（默认：https://api.moonshot.cn/v1/chat/completions）
- KIMI_MODEL: Kimi模型名称（默认：Kimi-K2）
- KIMI_TIMEOUT: API超时时间秒数（默认：60）
- LLM_TOP_N: LLM摘要TopN参数（默认：3）
- LLM_COVERAGE_THRESHOLD: 覆盖度阈值百分比（默认：80）
- LLM_FOCUS_THRESHOLD: 重点拦截阈值百分比（默认：10）

基本用法:
---------
# 位置参数方式（旧版兼容）
python Excel_Address_New_Modified.py <输入文件> [MTM表格] [输出目录] [开始日期] [结束日期]

# 命名参数方式（推荐）
python Excel_Address_New_Modified.py <输入文件> --start-date <开始日期> --end-date <结束日期>

参数说明:
---------
必需参数：
  输入文件                输入Excel文件路径（如：持续落入D等级 30天服务单明细.xlsx）

可选位置参数：
  MTM表格                 MTM映射表Excel路径（默认：mtm.xlsx）
  输出目录                输出目录路径（默认：output）
  开始日期                筛选开始日期（格式：YYYY-MM-DD 或 YYYY/MM/DD）
  结束日期                筛选结束日期（格式：YYYY-MM-DD 或 YYYY/MM/DD）

可选命名参数：
  --start-date            筛选开始日期（优先于位置参数）
  --end-date              筛选结束日期（优先于位置参数）
  --generate-ppt          生成PPT报告
  --use-llm               调用Kimi LLM生成详情页智能摘要（需配置KIMI_API_KEY）
  --ppt-template          PPT模板文件路径
  --ppt-path              输出PPT文件名（默认：output/report.pptx）
  --llm-timeout           LLM请求超时时间秒数（默认：60）
  --llm-top-n             LLM摘要TopN参数（默认：3）
  --llm-coverage          LLM摘要覆盖阈值百分比（默认：80）
  --llm-focus             LLM摘要重点拦截阈值百分比（默认：10）
  --skip-db               跳过数据库检查和导入
  --test-kimi             测试Kimi API连通性后退出

使用示例:
---------
# 示例1: 基本使用（仅生成Excel和图表）
python Excel_Address_New_Modified.py "持续落入D等级 30天服务单明细.xlsx"

# 示例2: 指定MTM映射表和输出目录
python Excel_Address_New_Modified.py "数据.xlsx" "mtm.xlsx" "output"

# 示例3: 指定日期范围
python Excel_Address_New_Modified.py "数据.xlsx" "mtm.xlsx" "output" "2025-07-01" "2025-07-18"

# 示例4: 使用命名参数指定日期范围
python Excel_Address_New_Modified.py "数据.xlsx" --start-date "2025-07-01" --end-date "2025-07-18"

# 示例5: 生成PPT报告（不使用LLM）
python Excel_Address_New_Modified.py "数据.xlsx" --generate-ppt

# 示例6: 生成PPT报告并使用LLM生成智能摘要
# 方式1: 直接在代码第171行配置API Key（推荐）
python Excel_Address_New_Modified.py "数据.xlsx" --generate-ppt --use-llm

# 方式2: 通过环境变量配置
export KIMI_API_KEY="your_api_key_here"  # Linux/Mac
set KIMI_API_KEY=your_api_key_here       # Windows
python Excel_Address_New_Modified.py "数据.xlsx" --generate-ppt --use-llm

# 示例7: 完整配置（自定义PPT路径、LLM参数）
python Excel_Address_New_Modified.py "数据.xlsx" \
    --start-date "2025-07-01" \
    --end-date "2025-07-18" \
    --generate-ppt \
    --use-llm \
    --ppt-path "我的报告.pptx" \
    --llm-timeout 90 \
    --llm-top-n 5

# 示例8: 跳过数据库检查（首次运行或无数据库时）
python Excel_Address_New_Modified.py "数据.xlsx" --skip-db --generate-ppt

# 示例9: 测试Kimi API连通性（配置API Key后首次使用时推荐）
python Excel_Address_New_Modified.py "数据.xlsx" --test-kimi

输出结果:
---------
output/
├── 审核原因统计.xlsx                    # 四种审核原因统计表
├── 7天无理由_机型分布.xlsx               # 7天无理由机型分布表
├── 非7天无理由_机型分布.xlsx             # 非7天无理由机型分布表
├── 审核原因占比.png                     # 审核原因饼图
├── 7天无理由_机型分布.png                # 7天无理由机型饼图
├── 非7天无理由_机型分布.png              # 非7天无理由机型饼图
├── 分析报告.txt                         # 文本分析报告
├── report.pptx                          # PPT报告（使用--generate-ppt时）
└── 详细数据/
    ├── 7天无理由/
    │   └── [机型名称]/
    │       ├── [机型]_7天无理由_分类频次.xlsx
    │       ├── [机型]_7天无理由_柱状图.png
    │       └── [机型]_7天无理由_详细数据.xlsx
    └── 非7天无理由/
        └── [机型名称]/
            ├── [机型]_非7天无理由_分类频次.xlsx
            ├── [机型]_非7天无理由_柱状图.png
            └── [机型]_非7天无理由_详细数据.xlsx

注意事项:
---------
1. Excel文件首列必须是日期列
2. 数据库配置默认为本地MySQL（localhost:3306/local_qcr）
3. 使用--use-llm前必须配置KIMI_API_KEY（在代码第171行或设置环境变量）
4. 生成的PPT采用空白布局，支持中文字体（微软雅黑）
5. 图片路径不存在时会跳过该图片，不影响其他内容生成
6. LLM生成失败时会自动降级为本地模板文本
7. 建议不要将包含真实API Key的代码提交到公开的版本控制系统

作者: KC
版本: 2.0
更新日期: 2025-10
=============================================================================
"""

import argparse
import json
import os
import re
import sys
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy import create_engine

# 设置matplotlib为非交互式后端（避免tkinter相关警告）
import matplotlib
matplotlib.use('Agg')  # 必须在导入pyplot之前设置
import matplotlib.pyplot as plt

# 设置中文字体（Windows 示例） 
matplotlib.rcParams['font.family'] = ['SimHei', 'Microsoft YaHei', 'DejaVu Sans']
matplotlib.rcParams['axes.unicode_minus'] = False


# 处理字体警告 - 使用更兼容的字体设置
import warnings
warnings.filterwarnings("ignore", category=UserWarning, message=".*Glyph.*missing.*")


# -----------------------------
# API密钥配置（请在此处配置您的Kimi API Key）
# -----------------------------
# 方式1: 直接在代码中配置（不推荐提交到版本控制）
DEFAULT_KIMI_API_KEY = "sk-z4mdCQLUIpPYoMwz7CMTonTHT8rgzgiaDOkkut5AJaHgU8wh"

# 方式2: 从环境变量读取（推荐，环境变量优先级高于代码配置）
# Windows: set KIMI_API_KEY=your_api_key
# Linux/Mac: export KIMI_API_KEY=your_api_key

# 最终使用的API Key（优先使用环境变量，其次使用代码配置）
KIMI_API_KEY = os.getenv("KIMI_API_KEY", DEFAULT_KIMI_API_KEY)

# -----------------------------
# 数据库配置
# -----------------------------
DB_CONFIG = {
    'host': 'localhost',
    'port': 3306,
    'user': 'root',
    'password': '0929',
    'database': 'local_qcr'
}

# -----------------------------
# 常量配置
# -----------------------------
DEFAULT_OUTPUT_DIR = "output"
DEFAULT_MTM_FILE = "mtm.xlsx"
DEFAULT_PPT_PATH = "report.pptx"
KIMI_API_URL = os.getenv("KIMI_API_URL", "https://api.moonshot.cn/v1/chat/completions")
KIMI_MODEL = os.getenv("KIMI_MODEL", "kimi-k2-0905-preview")
KIMI_TIMEOUT = int(os.getenv("KIMI_TIMEOUT", "60"))
LLM_TOP_N = int(os.getenv("LLM_TOP_N", "3"))
LLM_COVERAGE_THRESHOLD = float(os.getenv("LLM_COVERAGE_THRESHOLD", "80"))
LLM_FOCUS_THRESHOLD = float(os.getenv("LLM_FOCUS_THRESHOLD", "10"))

# PPT字体配置（可根据需要修改）
DEFAULT_PPT_FONT = "微软雅黑"  # 可选：宋体、黑体、Arial、SimHei等
PPT_TITLE_FONT_SIZE = 28       # 首页标题字号
PPT_SUBTITLE_FONT_SIZE = 28    # 详情页标题字号
PPT_BODY_FONT_SIZE = 14        # 正文字号


class LLMGenerationError(Exception):
    pass


# -----------------------------
# LLM相关辅助函数
# -----------------------------
def dataframe_to_category_rows(df: pd.DataFrame) -> List[Dict[str, str]]:
    rows = []
    for _, row in df.iterrows():
        rows.append({
            "Category": str(row.get("分类", "")),
            "Count": str(row.get("次数", "")),
            "Share": str(row.get("占比", ""))
        })
    return rows


def build_prompt_payload(category_rows: List[Dict[str, str]], top_n: int, coverage_threshold: float, focus_threshold: float) -> Dict[str, str]:
    table_lines = ["分类\t频次\t占比"]
    for row in category_rows:
        table_lines.append(f"{row['Category']}\t{row['Count']}\t{row['Share']}")
    table_text = "\n".join(table_lines)

    prompt = f"""# 角色
你是一名PC电脑制造业的质量管理专家与用户反馈分析专家。你的任务是在严格依赖输入表格（包含列：分类、频次、占比）的前提下，不引入外部信息、不自行计算/重算占比，输出高度凝练的核心观点与可执行建议，用于问题拦截与后续复现/根因分析。

## 输入数据表格
{table_text}

## 技能
### 技能 1: 生成核心观点
1. 输入包含"分类""频次""占比"的表格数据。
2. 开篇交代样本处理与总量（如"去除无效后共 N 项"），并直接点名 Top-N：
    - 采用紧凑体例："分类名*频次（占比）"（示例：无法开机*20（27.0%））。
3. 输出剩余分类情况："其余问题分布较为分散，无明显集中性"。
4. Top-N需要满足：
    - 必须输出Top1。
    - 对于Top2和Top3分别需要大于等于15%才可以被输出。
    - 对于Top4可以不输出。

### 技能 2: 生成可执行建议
1. 根据生成的核心观点中的Top-N分类。
2. 明确重点拦截清单，并给出下一步：
    - 给出拦截建议，一般拦截Top-N的机型，如果分类问题的频次较少可以不拦截，话术参考："建议对死机，无法开机等机器进行退机拦截处理，做进一步分析。"
    - 若分类名已明确指向方向（如"适配器-无法充电"），可给出极简方向性线索（从质量管理的角度，给出质量问题的探索方向，要求精简专业），避免越界推断。
    - 对于无理由退机的分类无需给出建议，直接忽略即可

## 限制:
- 不得计算/重算占比：不得基于频次推导占比，不得改写任何单项占比。
- 零幻觉：不添加输入表格之外的类别、原因或数据。
- 保留原词：引用分类名时保持与输入一致（除去多余空格）。
- 风格与数值：中文为主；百分比以输入为准，展示到2位小数（如输入非2位小数，原样输出或四舍五入但需注明）。
- 输出必须按照规定的格式和要求进行组织，不能偏离框架要求。
"""

    return {
        "role": "user",
        "content": prompt
    }


def call_kimi_api(messages: List[Dict[str, str]], timeout: int) -> str:
    api_key = KIMI_API_KEY
    if not api_key or api_key == "":
        raise LLMGenerationError("未配置KIMI_API_KEY，请在代码第171行或环境变量中配置API Key")

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    payload = {
        "model": KIMI_MODEL,
        "messages": messages,
        "temperature": 0.2
    }

    try:
        response = requests.post(
            KIMI_API_URL,
            headers=headers,
            data=json.dumps(payload),
            timeout=timeout
        )
    except requests.RequestException as exc:
        raise LLMGenerationError(f"Kimi API 请求异常: {exc}")

    if response.status_code != 200:
        raise LLMGenerationError(f"Kimi API 请求失败: {response.status_code} - {response.text}")

    data = response.json()
    try:
        return data["choices"][0]["message"]["content"].strip()
    except (KeyError, IndexError) as exc:
        raise LLMGenerationError(f"Kimi API 响应解析失败: {exc}")


def test_kimi_connection() -> bool:
    """
    测试Kimi API连通性
    发送一个简单的问候消息来验证API配置是否正确
    
    Returns:
        bool: 连接成功返回True，失败返回False
    """
    print("\n" + "="*60)
    print("🔍 Kimi API 连通性测试")
    print("="*60)
    
    # 检查API Key配置
    print(f"\n1. 检查API Key配置...")
    api_key = KIMI_API_KEY
    if not api_key or api_key == "":
        print("   ❌ 错误: 未配置KIMI_API_KEY")
        print("   请在代码第179行修改 DEFAULT_KIMI_API_KEY 或设置环境变量 KIMI_API_KEY")
        return False
    
    # 显示配置信息（隐藏部分密钥）
    masked_key = api_key[:10] + "..." + api_key[-8:] if len(api_key) > 18 else "***"
    print(f"   ✓ API Key: {masked_key}")
    print(f"   ✓ API URL: {KIMI_API_URL}")
    print(f"   ✓ 模型: {KIMI_MODEL}")
    
    # 发送测试请求
    print(f"\n2. 发送测试请求...")
    test_message = {
        "role": "user",
        "content": "你好，请简单回复'连接成功'即可。"
    }
    
    try:
        response = call_kimi_api([test_message], timeout=30)
        print(f"   ✓ 请求成功!")
        print(f"\n3. Kimi 响应:")
        print(f"   {response}")
        
        print("\n" + "="*60)
        print("✅ Kimi API 连接测试成功！")
        print("="*60 + "\n")
        return True
        
    except LLMGenerationError as exc:
        print(f"   ❌ 请求失败: {exc}")
        print("\n" + "="*60)
        print("❌ Kimi API 连接测试失败")
        print("="*60)
        print("\n可能的原因:")
        print("1. API Key 不正确或已过期")
        print("2. 网络连接问题")
        print("3. API 服务暂时不可用")
        print("4. API URL 或模型名称配置错误")
        print("\n请检查配置后重试。\n")
        return False


def generate_llm_summary(category_df: pd.DataFrame, timeout: int, top_n: int, coverage_threshold: float, focus_threshold: float) -> str:
    category_rows = dataframe_to_category_rows(category_df)
    if not category_rows:
        raise LLMGenerationError("分类数据为空，无法生成LLM摘要")

    message = build_prompt_payload(category_rows, top_n, coverage_threshold, focus_threshold)
    return call_kimi_api([message], timeout)


def default_llm_fallback(clean_model: str, suffix: str, total_records: int) -> str:
    return (
        "核心观点（Human-Readable Core Insights）\n"
        f"- 样本：{clean_model}{suffix}共 {total_records} 条，暂未能生成自动化摘要。\n"
        "- 暂未获取模型结论，建议人工复核分类表。"
    )


# -----------------------------
# PPT生成相关函数
# -----------------------------
def add_textbox_with_content(slide, left, top, width, height, text, font_name=None, font_size=None, bold=False):
    """
    在幻灯片中添加文本框
    
    Args:
        slide: 幻灯片对象
        left, top, width, height: 文本框位置和尺寸（Inches对象）
        text: 文本内容
        font_name: 字体名称（默认使用DEFAULT_PPT_FONT）
        font_size: 字号（默认使用PPT_BODY_FONT_SIZE）
        bold: 是否加粗
    """
    if font_name is None:
        font_name = DEFAULT_PPT_FONT
    if font_size is None:
        font_size = PPT_BODY_FONT_SIZE
        
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    return textbox


def add_image(slide, image_path: str, left, top, width=None, height=None):
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def build_homepage_slide(prs: Presentation, payload: Dict[str, Any]):
    """
    生成首页幻灯片
    布局：标题（顶部）→ 正文（中间）→ 图表（底部横向排列）
    """
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # ========== 标题区域（顶部）==========
    title_text = "落入D等级 数据汇总分析结果"
    add_textbox_with_content(
        slide, 
        Inches(0.5), Inches(0.3), 
        Inches(9), Inches(0.8), 
        title_text, 
        font_size=PPT_TITLE_FONT_SIZE, 
        bold=True
    )

    # ========== 正文区域（中间）==========
    bullet_left = Inches(0.5)
    bullet_top = Inches(1.2)
    bullet_box = slide.shapes.add_textbox(bullet_left, bullet_top, Inches(9), Inches(2.5))
    frame = bullet_box.text_frame
    frame.word_wrap = True
    frame.clear()

    # 提取数据
    start_str, end_str = payload.get("coverage_period", ("-", "-"))
    week_start, week_end = payload.get("week_range", ("-", "-"))
    product_list = payload.get("unique_models", [])
    product_text = "、".join(product_list[:4]) if product_list else "暂无"
    total_records = payload.get("total_records", 0)

    # Bullet 1: 基本信息
    p1 = frame.add_paragraph()
    p1.text = (
        f"{week_start}-{week_end}共收到落入D等级产品数据{total_records}条，覆盖周期为{start_str}-{end_str}，"
        f"产品为{product_text}等，共计{len(product_list)}款。"
    )
    p1.level = 0
    p1.font.name = DEFAULT_PPT_FONT
    p1.font.size = Pt(PPT_BODY_FONT_SIZE)

    # Bullet 2: 审核原因占比
    reason_df = payload.get("reason_stats", pd.DataFrame())
    if not reason_df.empty:
        row_map = {row["审核原因"]: f"{row['占比']}%" for _, row in reason_df.iterrows()}
        p2 = frame.add_paragraph()
        p2.text = (
            f"审核原因中，7天无理由占比{row_map.get('7天无理由', '0%')}，15天质量换新占比{row_map.get('15天质量换新', '0%')}，"
            f"质量维修占比{row_map.get('质量维修', '0%')}，180天只换不修占比{row_map.get('180天只换不修', '0%')}。"
        )
        p2.level = 0
        p2.font.name = DEFAULT_PPT_FONT
        p2.font.size = Pt(PPT_BODY_FONT_SIZE)

    # Bullet 3: 七天无理由机型
    model_7d_df = payload.get("model_7d_dist", pd.DataFrame())
    if not model_7d_df.empty:
        top_items = model_7d_df.head(4)
        parts = [f"{row['机型名称']}占比{row['占比']}%" for _, row in top_items.iterrows()]
        p3 = frame.add_paragraph()
        p3.text = "七天无理由中，" + "，".join(parts) + "。"
        p3.level = 0
        p3.font.name = DEFAULT_PPT_FONT
        p3.font.size = Pt(PPT_BODY_FONT_SIZE)

    # Bullet 4: 非七天无理由机型
    model_non7d_df = payload.get("model_non_7d_dist", pd.DataFrame())
    if not model_non7d_df.empty:
        top_items = model_non7d_df.head(4)
        parts = [f"{row['机型名称']}占比{row['占比']}%" for _, row in top_items.iterrows()]
        p4 = frame.add_paragraph()
        p4.text = "非七天无理由中，" + "，".join(parts) + "。"
        p4.level = 0
        p4.font.name = DEFAULT_PPT_FONT
        p4.font.size = Pt(PPT_BODY_FONT_SIZE)

    # ========== 图表区域（底部横向排列）==========
    # 3张饼图横向排列，Y轴统一为4.2英寸
    chart_y = Inches(4.2)
    chart_width = Inches(2.5)
    
    # 图1: 审核原因饼图（左）
    add_image(slide, payload.get("reason_chart_path"), Inches(1.0), chart_y, width=chart_width)
    
    # 图2: 7天机型饼图（中）
    add_image(slide, payload.get("model_7d_chart_path"), Inches(4.2), chart_y, width=chart_width)
    
    # 图3: 非7天机型饼图（右）
    add_image(slide, payload.get("model_non7d_chart_path"), Inches(7.4), chart_y, width=chart_width)


def build_detail_slide(prs: Presentation, model_name: str, suffix: str, entry: Dict[str, Any], use_llm: bool, llm_params: Dict[str, Any]):
    """
    生成详情页幻灯片
    布局：标题（顶部）→ 正文（中间）→ 图表（底部居中）
    """
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # ========== 标题区域（顶部）==========
    title = f"{model_name} {suffix}分类"
    add_textbox_with_content(
        slide, 
        Inches(0.5), Inches(0.3), 
        Inches(9), Inches(0.8), 
        title, 
        font_size=PPT_SUBTITLE_FONT_SIZE, 
        bold=True
    )

    # ========== 正文区域（中间）==========
    # 生成LLM内容或使用fallback，并打印进度
    print(f"→ 开始生成详情页：机型='{model_name}', 类型='{suffix}'")
    text_content = ""
    if use_llm:
        try:
            print(f"   调用Kimi生成观点中... 机型='{model_name}', 类型='{suffix}'")
            text_content = generate_llm_summary(
                entry.get("category_df", pd.DataFrame()),
                timeout=llm_params["timeout"],
                top_n=llm_params["top_n"],
                coverage_threshold=llm_params["coverage"],
                focus_threshold=llm_params["focus"]
            )
            print(f"   ✓ Kimi返回观点（{model_name}-{suffix}）：\n{text_content}\n")
        except LLMGenerationError as exc:
            print(f"   ✗ LLM生成失败[{model_name}-{suffix}]: {exc}")
            text_content = default_llm_fallback(entry.get("clean_model", ""), suffix, entry.get("total_records", 0))
            print(f"   → 使用本地模板观点（{model_name}-{suffix}）：\n{text_content}\n")
    else:
        text_content = default_llm_fallback(entry.get("clean_model", ""), suffix, entry.get("total_records", 0))
        print(f"   （未启用LLM）本地模板观点（{model_name}-{suffix}）：\n{text_content}\n")

    # 添加正文文本框（中间区域）
    add_textbox_with_content(
        slide, 
        Inches(0.5), Inches(1.2), 
        Inches(9), Inches(2.8), 
        text_content,
        font_size=PPT_BODY_FONT_SIZE
    )

    # ========== 图表区域（底部居中）==========
    # 柱状图居中显示
    chart_width = Inches(4.5)
    chart_left = Inches(2.75)  # (10 - 4.5) / 2 = 2.75 实现居中
    chart_top = Inches(4.5)
    
    add_image(slide, entry.get("chart_path"), chart_left, chart_top, width=chart_width)
    print(f"← 完成：机型='{model_name}', 类型='{suffix}' 的详情页\n")


def generate_ppt(summary_payload: Dict[str, Any], output_path: Path, template_path: Optional[Path], use_llm: bool, llm_params: Dict[str, Any]):
    if template_path and template_path.exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    print("\n===== 开始生成PPT =====")
    print("生成首页...")
    build_homepage_slide(prs, summary_payload)

    model_details = summary_payload.get("model_details", {})
    for model_name, entries in model_details.items():
        for suffix in ["7天无理由", "非7天无理由"]:
            entry = entries.get(suffix)
            if not entry:
                continue
            print(f"准备生成：机型='{model_name}', 类型='{suffix}' 的详情页...")
            build_detail_slide(prs, model_name, suffix, entry, use_llm, llm_params)

    prs.save(str(output_path))
    print(f"PPT报告已生成：{output_path}")
    print("===== 生成完成 =====\n")


def parse_arguments():
    parser = argparse.ArgumentParser(
        description="处理QCR数据并生成分析报告、图表及可选PPT"
    )
    parser.add_argument("input_file", help="输入Excel文件路径")
    parser.add_argument("mtm_file", nargs="?", default=DEFAULT_MTM_FILE, help="MTM映射表Excel路径")
    parser.add_argument("output_dir", nargs="?", default=DEFAULT_OUTPUT_DIR, help="输出目录路径")
    parser.add_argument("start_date_arg", nargs="?", default=None, help="开始日期 (YYYY-MM-DD)")
    parser.add_argument("end_date_arg", nargs="?", default=None, help="结束日期 (YYYY-MM-DD)")

    parser.add_argument("--start-date", dest="start_date_opt", help="开始日期 (YYYY-MM-DD)")
    parser.add_argument("--end-date", dest="end_date_opt", help="结束日期 (YYYY-MM-DD)")
    parser.add_argument("--generate-ppt", action="store_true", help="生成PPT报告")
    parser.add_argument("--use-llm", action="store_true", help="调用Kimi LLM生成详情页摘要")
    parser.add_argument("--ppt-template", dest="ppt_template", default=None, help="PPT模板文件路径")
    parser.add_argument("--ppt-path", dest="ppt_path", default=None, help="输出PPT文件名")
    parser.add_argument("--llm-timeout", dest="llm_timeout", type=int, default=KIMI_TIMEOUT, help="LLM请求超时时间(秒)")
    parser.add_argument("--llm-top-n", dest="llm_top_n", type=int, default=LLM_TOP_N, help="LLM摘要TopN参数")
    parser.add_argument(
        "--llm-coverage", dest="llm_coverage", type=float, default=LLM_COVERAGE_THRESHOLD,
        help="LLM摘要覆盖阈值(%)"
    )
    parser.add_argument(
        "--llm-focus", dest="llm_focus", type=float, default=LLM_FOCUS_THRESHOLD,
        help="LLM摘要重点拦截阈值(%)"
    )
    parser.add_argument("--skip-db", action="store_true", help="跳过数据库检查和导入")
    parser.add_argument("--test-kimi", action="store_true", help="测试Kimi API连通性后退出")

    return parser.parse_args()

# -----------------------------
# 工具函数：清理文件名中的非法字符
# -----------------------------
def sanitize_filename(filename):
    """清理文件名中的非法字符"""
    # Windows非法字符：<>:"/\|?*
    illegal_chars = r'[<>:\"/\\|?*]'
    # 替换为空格
    filename = re.sub(illegal_chars, ' ', filename)
    # 去除前后空格
    filename = filename.strip()
    # 限制长度
    if len(filename) > 200:
        filename = filename[:200]
    return filename

# -----------------------------
# 数据库工具函数
# -----------------------------
def check_and_import_new_data(df):
    """检查数据库中不存在的服务单号并导入新数据"""
    try:
        print("开始连接数据库...")
        # 创建数据库连接
        connection_string = (
            f"mysql+pymysql://{DB_CONFIG['user']}:{DB_CONFIG['password']}@"
            f"{DB_CONFIG['host']}:{DB_CONFIG['port']}/{DB_CONFIG['database']}"
        )
        engine = create_engine(connection_string)
        
        # 检查数据框中是否有服务单号列
        service_order_column = None
        for col in df.columns:
            if '服务单号' in str(col):
                service_order_column = col
                break
        
        if service_order_column is None:
            print("警告：未找到'服务单号'列，跳过数据库检查")
            return df
        
        # 获取当前数据中的服务单号
        current_service_orders = df[service_order_column].dropna().astype(str).tolist()
        print(f"当前数据包含 {len(current_service_orders)} 个服务单号")
        
        # 查询数据库中已存在的服务单号
        try:
            # 先检查表是否存在
            table_exists = pd.read_sql(
                "SELECT COUNT(*) as count FROM information_schema.tables WHERE table_schema = %s AND table_name = 'qcr_data'", 
                engine, 
                params=(DB_CONFIG['database'],)
            )['count'].iloc[0] > 0
            
            if table_exists:
                existing_service_orders = pd.read_sql(
                    "SELECT service_order_id FROM qcr_data", 
                    engine
                )['service_order_id'].astype(str).tolist()
                print(f"数据库中已存在 {len(existing_service_orders)} 个服务单号")
            else:
                print("数据库表qcr_data不存在，将创建新表")
                existing_service_orders = []
        except Exception as e:
            print(f"查询数据库失败，假设数据库为空: {e}")
            existing_service_orders = []
        
        # 筛选出数据库中不存在的新服务单号
        new_service_orders = [
            order for order in current_service_orders 
            if order not in existing_service_orders
        ]
        print(f"新服务单号数量: {len(new_service_orders)}")
        
        # 筛选新数据
        df_new = df[df[service_order_column].astype(str).isin(new_service_orders)].copy()
        
        if len(df_new) == 0:
            print("没有新数据需要导入和分析")
            return df_new
        
        # 准备导入数据库的数据
        df_to_import = df_new.copy()
        
        # 重命名列以匹配数据库字段
        column_mapping = {}
        for col in df_to_import.columns:
            if '服务单号' in str(col):
                column_mapping[col] = 'service_order_id'
            elif '日期' in str(col):
                column_mapping[col] = 'date'
            elif '订单号' in str(col):
                column_mapping[col] = 'order_id'
            elif '问题描述' in str(col):
                column_mapping[col] = 'issue_description'
            elif 'SKU' in str(col):
                column_mapping[col] = 'sku'
            elif 'SN编码' in str(col):
                column_mapping[col] = 'sn_code'
            elif '客户账号' in str(col) or '客户账户' in str(col):
                column_mapping[col] = 'customer_account'
            elif '商品名称' in str(col):
                column_mapping[col] = 'product_name'
            elif 'MTM' in str(col):
                column_mapping[col] = 'mtm'
            elif '审核原因' in str(col):
                column_mapping[col] = 'audit_reason'
            elif '问题分类' in str(col) and '一' not in str(col):
                column_mapping[col] = 'issue_category'
            elif '分类' in str(col) and '问题' not in str(col):
                column_mapping[col] = 'category'
        
        # 应用列映射
        df_to_import = df_to_import.rename(columns=column_mapping)
        
        # 确保必需的列存在
        required_db_columns = [
            'service_order_id', 'date', 'order_id', 'issue_description', 
            'sku', 'sn_code', 'customer_account', 'product_name', 
            'mtm', 'audit_reason', 'issue_category', 'category'
        ]
        
        for col in required_db_columns:
            if col not in df_to_import.columns:
                df_to_import[col] = ''
        
        # 数据类型转换和清洗
        # 处理日期
        if 'date' in df_to_import.columns:
            df_to_import['date'] = pd.to_datetime(df_to_import['date'], errors='coerce').dt.strftime('%Y-%m-%d')
            df_to_import['date'] = df_to_import['date'].fillna('1900-01-01')
        
        # 处理数值列
        numeric_columns = ['service_order_id', 'order_id', 'sku']
        for col in numeric_columns:
            if col in df_to_import.columns:
                df_to_import[col] = pd.to_numeric(df_to_import[col], errors='coerce').astype('Int64')
        
        # 处理字符串列
        string_columns = [
            'issue_description', 'sn_code', 'customer_account', 
            'product_name', 'mtm', 'audit_reason', 
            'issue_category', 'category'
        ]
        for col in string_columns:
            if col in df_to_import.columns:
                df_to_import[col] = df_to_import[col].fillna('').astype(str).str.strip()
                # 字符串长度限制
                max_lengths = {
                    'issue_description': 500,
                    'sn_code': 100,
                    'customer_account': 100,
                    'product_name': 200,
                    'mtm': 50,
                    'audit_reason': 100,
                    'issue_category': 100,
                    'category': 100
                }
                if col in max_lengths:
                    df_to_import[col] = df_to_import[col].str[:max_lengths[col]]
        
        # 删除必需字段为空的行
        if 'service_order_id' in df_to_import.columns:
            df_to_import = df_to_import.dropna(subset=['service_order_id'])
        
        # 只选择数据库需要的列
        df_to_import = df_to_import[required_db_columns]
        
        # 导入数据到数据库
        if len(df_to_import) > 0:
            try:
                df_to_import.to_sql(
                    'qcr_data', 
                    engine, 
                    if_exists='append', 
                    index=False,
                    method='multi'
                )
                print(f"成功导入 {len(df_to_import)} 条新记录到数据库")
            except Exception as e:
                print(f"导入数据到数据库失败: {e}")
                print("将继续分析当前数据，但新数据不会保存到数据库")
        
        # 返回新数据用于后续分析
        return df_new
        
    except Exception as e:
        print(f"数据库操作失败: {e}")
        print("将继续分析原始数据，跳过数据库检查和导入")
        return df

# -----------------------------
# 1. 命令行参数处理
# -----------------------------
args = parse_arguments()

# -----------------------------
# Kimi API 连通性测试（如果启用）
# -----------------------------
if args.test_kimi:
    test_kimi_connection()
    sys.exit(0)  # 测试完成后退出

file_path = Path(args.input_file)
if not file_path.exists():
    print(f"错误: 文件 '{file_path}' 不存在")
    sys.exit(1)

mtm_file_path = Path(args.mtm_file)
if not mtm_file_path.exists():
    print(f"警告: MTM表格文件 '{mtm_file_path}' 不存在，将使用原始MTM值")
    use_mtm_mapping = False
else:
    use_mtm_mapping = True

out_dir = Path(args.output_dir)
out_dir.mkdir(parents=True, exist_ok=True)
sheet_name = 0  # 默认第一张表

# -----------------------------
# 日期解析函数
# -----------------------------
def parse_date(date_str):
    """尝试解析多种日期格式"""
    for fmt in ("%Y-%m-%d", "%Y/%m/%d"):
        try:
            return datetime.strptime(date_str, fmt).date()
        except ValueError:
            continue
    raise ValueError(f"无法解析日期: {date_str}，请使用 YYYY-MM-DD 或 YYYY/MM/DD 格式")


def format_percentage(value: float, decimals: int = 1) -> str:
    return f"{value:.{decimals}f}%"


def parse_percentage(value) -> Optional[float]:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    value_str = str(value).strip().replace('%', '')
    if not value_str:
        return None
    try:
        return float(value_str)
    except ValueError:
        return None


def get_week_workday_range(reference_date: Optional[date] = None) -> Tuple[str, str]:
    today = reference_date if reference_date else datetime.today().date()
    monday = today - timedelta(days=today.weekday())
    friday = monday + timedelta(days=4)
    return monday.strftime("%Y/%m/%d"), friday.strftime("%Y/%m/%d")


def determine_coverage_range(df: pd.DataFrame, date_column: str, start_date: Optional[date], end_date: Optional[date]) -> Tuple[str, str]:
    if df.empty:
        return ("-", "-")
    
    # 确保日期类型一致
    actual_start = start_date if start_date else df[date_column].min()
    actual_end = end_date if end_date else df[date_column].max()
    
    # 将 pandas Timestamp 转换为 date 对象
    if hasattr(actual_start, 'date') and callable(actual_start.date):
        actual_start = actual_start.date()
    if hasattr(actual_end, 'date') and callable(actual_end.date):
        actual_end = actual_end.date()
    
    return actual_start.strftime("%Y/%m/%d"), actual_end.strftime("%Y/%m/%d")


def join_top_items(items: List[str], limit: int) -> str:
    filtered = [str(item) for item in items if str(item)]
    return "、".join(filtered[:limit])

# -----------------------------
# 获取日期范围参数
# -----------------------------
start_date = None
end_date = None

def resolve_date(raw_date):
    if not raw_date:
        return None
    try:
        return parse_date(raw_date)
    except ValueError as exc:
        print(f"警告: {exc}，将处理所有数据")
        return None

start_date = resolve_date(args.start_date_opt or args.start_date_arg)
end_date = resolve_date(args.end_date_opt or args.end_date_arg)

# -----------------------------
# 2. 读取数据
# -----------------------------
df = pd.read_excel(file_path, sheet_name=sheet_name)

# 假设第一列是日期列，转换为日期格式
date_column = df.columns[0]
df[date_column] = pd.to_datetime(df[date_column]).dt.date

# 根据日期范围筛选数据
if start_date and end_date:
    mask = (df[date_column] >= start_date) & (df[date_column] <= end_date)
    df = df[mask]
    print(f"已筛选 {start_date} 到 {end_date} 的数据，共 {len(df)} 条记录")
elif start_date:
    mask = df[date_column] >= start_date
    df = df[mask]
    print(f"已筛选 {start_date} 之后的数据，共 {len(df)} 条记录")
elif end_date:
    mask = df[date_column] <= end_date
    df = df[mask]
    print(f"已筛选 {end_date} 之前的数据，共 {len(df)} 条记录")

# -----------------------------
# 数据库检查和导入新数据
# -----------------------------
if args.skip_db:
    print("已跳过数据库检查与导入")
else:
    print("\n开始检查数据库中已存在的服务单号...")
    df = check_and_import_new_data(df)
    print(f"数据库检查后，剩余 {len(df)} 条新记录需要分析\n")

# -----------------------------
# 3. 去重MTM
# -----------------------------
# original_count = len(df)
# df = df.drop_duplicates(subset=['MTM'])
# print(f"已去重MTM，从 {original_count} 条记录减少到 {len(df)} 条记录")

# -----------------------------
# 4. 读取MTM映射表
# -----------------------------
if use_mtm_mapping:
    mtm_df = pd.read_excel(mtm_file_path, sheet_name=sheet_name, header=None)
    mtm_df.columns = ['MTM', '机型名称']
    mtm_mapping = dict(zip(mtm_df['MTM'], mtm_df['机型名称']))
    
    # 映射MTM到机型名称
    df['机型名称'] = df['MTM'].map(mtm_mapping).fillna(df['MTM'])
else:
    # 如果没有MTM映射表，使用原始MTM值作为机型名称
    df['机型名称'] = df['MTM']

# -----------------------------
# 5. 预计算常用条件
# -----------------------------
cond_7d = df["审核原因"] == "7天无理由"
cond_non_7d = df["审核原因"].isin(["15天质量换新", "180天只换不修", "质量维修"])

# 缓存中间结果
df_7d = df[cond_7d].copy()
df_non_7d = df[cond_non_7d].copy()

# 创建文件夹结构
detailed_dir_7d = out_dir / "详细数据" / "7天无理由"
detailed_dir_non7d = out_dir / "详细数据" / "非7天无理由"
detailed_dir_7d.mkdir(parents=True, exist_ok=True)
detailed_dir_non7d.mkdir(parents=True, exist_ok=True)

# -----------------------------
# 6. 统计四种审核原因
# -----------------------------
reasons = ["15天质量换新", "180天只换不修", "7天无理由", "质量维修"]
counts = {r: int((df["审核原因"] == r).sum()) for r in reasons}

summary_df = pd.DataFrame(list(counts.items()), columns=["审核原因", "数量"])
total_reason_count = summary_df["数量"].sum()
summary_df["占比"] = (summary_df["数量"] / total_reason_count * 100).round(2)
summary_df.to_excel(out_dir / "审核原因统计.xlsx", index=False)

plt.figure(figsize=(6, 6))
plt.pie(summary_df["数量"], labels=summary_df["审核原因"], autopct="%1.1f%%")
plt.title("审核原因占比")
plt.tight_layout()
reason_chart_path = out_dir / "审核原因占比.png"
plt.savefig(reason_chart_path)
plt.close()

# -----------------------------
# 7. 7天无理由机型分布
# -----------------------------
model_7d_dist = pd.DataFrame()
model_7d_chart_path = None
if len(df_7d) > 0:
    model_7d_dist = (
        df_7d["机型名称"]
        .value_counts()
        .rename_axis("机型名称")
        .reset_index(name="数量")
        .assign(占比=lambda x: (x["数量"] / x["数量"].sum() * 100).round(1))
    )
    model_7d_dist.to_excel(out_dir / "7天无理由_机型分布.xlsx", index=False)

    plt.figure(figsize=(8, 8))
    plt.pie(model_7d_dist["数量"], labels=model_7d_dist["机型名称"], autopct="%1.1f%%")
    plt.title("7天无理由 - 机型分布")
    plt.tight_layout()
    model_7d_chart_path = out_dir / "7天无理由_机型分布.png"
    plt.savefig(model_7d_chart_path)
    plt.close()
else:
    print("警告：7天无理由数据为空")

# -----------------------------
# 8. 非7天无理由机型分布
# -----------------------------
model_non_7d_dist = pd.DataFrame()
model_non_7d_chart_path = None
if len(df_non_7d) > 0:
    model_non_7d_dist = (
        df_non_7d["机型名称"]
        .value_counts()
        .rename_axis("机型名称")
        .reset_index(name="数量")
        .assign(占比=lambda x: (x["数量"] / x["数量"].sum() * 100).round(1))
    )
    model_non_7d_dist.to_excel(out_dir / "非7天无理由_机型分布.xlsx", index=False)

    plt.figure(figsize=(8, 8))
    plt.pie(model_non_7d_dist["数量"], labels=model_non_7d_dist["机型名称"], autopct="%1.1f%%")
    plt.title("非7天无理由 - 机型分布")
    plt.tight_layout()
    model_non_7d_chart_path = out_dir / "非7天无理由_机型分布.png"
    plt.savefig(model_non_7d_chart_path)
    plt.close()
else:
    print("警告：非7天无理由数据为空")

# -----------------------------
# 9. 按机型统计分类描述词频次
# -----------------------------
def build_model_issue_table(df_sub, suffix, detailed_dir):
    """为每个机型计算'分类'描述词频次，输出excel和柱状图，并生成详细数据文件"""
    if len(df_sub) == 0:
        print(f"警告：{suffix}数据为空，跳过机型分析")
        return []

    summaries = []
        
    # 非7天无理由数据：过滤掉问题描述为空的行
    if suffix == "非7天无理由":
        # 检查问题描述列是否存在
        if "问题描述" in df_sub.columns:
            # 过滤掉问题描述为空的行
            df_sub = df_sub[df_sub["问题描述"].notna() & (df_sub["问题描述"] != "")]
            print(f"已过滤空问题描述行，剩余 {len(df_sub)} 条记录")
        else:
            print("警告：未找到'问题描述'列，无法过滤空值")
    
    for model in df_sub["机型名称"].unique():
        # 清理机型名称用于文件夹和文件名
        clean_model = sanitize_filename(str(model))
        
        # 创建机型文件夹
        model_dir = detailed_dir / clean_model
        model_dir.mkdir(parents=True, exist_ok=True)
        
        # 获取该机型的所有数据
        model_data = df_sub[df_sub["机型名称"] == model].copy()
        
        # 统计分类频次
        sub = (
            model_data["分类"]
            .value_counts()
            .rename_axis("分类")
            .reset_index(name="次数")
        )

        if "次数" in sub.columns and sub["次数"].sum() > 0:
            sub["占比"] = (sub["次数"] / sub["次数"].sum() * 100).round(1)
        else:
            sub["占比"] = 0
        
        # 保存频次统计
        freq_filename = f"{clean_model}_{suffix}_分类频次.xlsx"
        freq_path = model_dir / freq_filename
        sub.to_excel(freq_path, index=False)

        # 为每个机型的所有分类生成一个综合详细数据文件
        detailed_filename = f"{clean_model}_{suffix}_详细数据.xlsx"
        detailed_path = model_dir / detailed_filename
        model_data.to_excel(detailed_path, index=False)
        
        # 生成柱状图
        plt.figure(figsize=(12, 6))
        bars = plt.bar(sub["分类"], sub["次数"])
        plt.xticks(rotation=45, ha="right")
        plt.title(f"{model} - {suffix} - 分类频次")
        
        # 添加数量标签
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height,
                    f'{int(height)}', ha='center', va='bottom')
        
        plt.tight_layout()
        
        chart_filename = f"{clean_model}_{suffix}_柱状图.png"
        chart_path = model_dir / chart_filename
        plt.savefig(chart_path)
        plt.close()
        
        print(f"已生成 {model} 的 {suffix} 数据，共 {len(sub)} 个分类，{len(model_data)} 条记录")

        model_summary = {
            "model": model,
            "clean_model": clean_model,
            "suffix": suffix,
            "category_df": sub,
            "chart_path": str(chart_path),
            "total_records": len(model_data)
        }
        summaries.append(model_summary)

    return summaries


# -----------------------------
# 修复 generate_analysis_report 函数中的列名问题
# -----------------------------
def generate_analysis_report(df, df_7d, df_non_7d, out_dir, start_date, end_date):
    """生成分析报告并保存到文本文件"""
    report_lines = []

    # 1. 落入D等级产品数据统计
    d_grade_data = df[df['审核原因'] == 'D等级']
    d_grade_count = len(d_grade_data)
    d_grade_models = d_grade_data['机型名称'].unique()
    d_grade_model_count = len(d_grade_models)
    report_lines.append(f"落入D等级产品数据：{d_grade_count} 条")
    report_lines.append(f"覆盖周期：{start_date} 至 {end_date}")
    report_lines.append(f"涉及机型：{', '.join(d_grade_models)}")
    report_lines.append(f"共计机型数量：{d_grade_model_count} 款\n")

    # 2. 审核原因占比
    reasons = ["7天无理由", "15天质量换新", "质量维修", "180天只换不修"]
    total_count = len(df)
    for reason in reasons:
        count = (df['审核原因'] == reason).sum()
        percentage = (count / total_count * 100) if total_count > 0 else 0
        report_lines.append(f"审核原因 - {reason}：{count} 条，占比 {percentage:.2f}%")
    report_lines.append("")

    # 3. 七天无理由机型占比
    if len(df_7d) > 0:
        model_7d_dist = (
            df_7d['机型名称']
            .value_counts()
            .rename_axis('机型名称')
            .reset_index(name='数量')
            .assign(占比=lambda x: (x['数量'] / x['数量'].sum() * 100).round(2))
        )
        report_lines.append("七天无理由机型占比：")
        for _, row in model_7d_dist.iterrows():
            report_lines.append(f"  {row['机型名称']}: {row['数量']} 条，占比 {row['占比']}%")
        report_lines.append("")

    # 4. 非七天无理由机型占比
    if len(df_non_7d) > 0:
        model_non_7d_dist = (
            df_non_7d['机型名称']
            .value_counts()
            .rename_axis('机型名称')
            .reset_index(name='数量')
            .assign(占比=lambda x: (x['数量'] / x['数量'].sum() * 100).round(2))
        )
        report_lines.append("非七天无理由机型占比：")
        for _, row in model_non_7d_dist.iterrows():
            report_lines.append(f"  {row['机型名称']}: {row['数量']} 条，占比 {row['占比']}%")
        report_lines.append("")

    # 5. 每个机型七天无理由的分类数据分析
    report_lines.append("每个机型七天无理由的分类数据分析：")
    for model in df_7d['机型名称'].unique():
        model_data = df_7d[df_7d['机型名称'] == model]
        total_comments = len(model_data)
        no_reason_count = (model_data['分类'] == '无理由退货').sum()
        no_reason_percentage = (no_reason_count / total_comments * 100) if total_comments > 0 else 0
        top_issues = (
            model_data['分类']
            .value_counts()
            .reset_index(name='次数')
            .rename(columns={'index': '分类'})  # 确保列名正确
        )

        # 调试信息：打印 top_issues 列名和数据
        print("七天无理由 - Top Issues:")
        print(top_issues.head())

        top_issues = top_issues[top_issues['次数'] >= 2].head(2)
        report_lines.append(f"  {model}:")
        report_lines.append(f"    评论总数：{total_comments}")
        report_lines.append(f"    无理由退货：{no_reason_count} 条，占比 {no_reason_percentage:.2f}%")
        for _, row in top_issues.iterrows():
            issue_percentage = (row['次数'] / total_comments * 100) if total_comments > 0 else 0
            report_lines.append(f"    Top问题：{row['分类']}，次数：{row['次数']}，占比：{issue_percentage:.2f}%")
    report_lines.append("")

    # 6. 每个机型非七天无理由的分类数据分析
    report_lines.append("每个机型非七天无理由的分类数据分析：")
    for model in df_non_7d['机型名称'].unique():
        model_data = df_non_7d[df_non_7d['机型名称'] == model]
        total_comments = len(model_data)
        top_issues = (
            model_data['分类']
            .value_counts()
            .reset_index(name='次数')
            .rename(columns={'index': '分类'})  # 确保列名正确
        )

        # 调试信息：打印 top_issues 列名和数据
        print("非七天无理由 - Top Issues:")
        print(top_issues.head())

        top_issues = top_issues[top_issues['次数'] >= 2].head(2)
        report_lines.append(f"  {model}:")
        report_lines.append(f"    有效评论总数：{total_comments}")
        for _, row in top_issues.iterrows():
            issue_percentage = (row['次数'] / total_comments * 100) if total_comments > 0 else 0
            report_lines.append(f"    Top问题：{row['分类']}，次数：{row['次数']}，占比：{issue_percentage:.2f}%")
    report_lines.append("")

    # 7. 总结
    report_lines.append("总结：")
    report_lines.append(f"本次报告时间覆盖：{start_date} 至 {end_date}")
    report_lines.append(f"覆盖机型：{', '.join(df['机型名称'].unique())}")
    report_lines.append("非七天无理由分类中，以下机型的问题较为突出：")
    for model in df_non_7d['机型名称'].unique():
        model_data = df_non_7d[df_non_7d['机型名称'] == model]
        top_issues = (
            model_data['分类']
            .value_counts()
            .reset_index(name='次数')
            .rename(columns={'index': '分类'})  # 确保列名正确
        )
        top_issues = top_issues[top_issues['次数'] >= 2].head(2)
        for _, row in top_issues.iterrows():
            report_lines.append(f"  {model} - {row['分类']}：{row['次数']} 次")

    # 保存报告到文件
    report_path = out_dir / "分析报告.txt"
    with open(report_path, "w", encoding="utf-8") as f:
        f.write("\n".join(report_lines))

    print(f"分析报告已生成：{report_path}")

summary_payload = {
    "start_date": start_date,
    "end_date": end_date,
    "week_range": get_week_workday_range(),
    "coverage_period": determine_coverage_range(df, date_column, start_date, end_date),
    "total_records": len(df),
    "unique_models": sorted(df["机型名称"].dropna().unique().tolist()),
    "reason_stats": summary_df,
    "model_7d_dist": model_7d_dist,
    "model_non_7d_dist": model_non_7d_dist,
    "reason_chart_path": str(reason_chart_path),
    "model_7d_chart_path": str(model_7d_chart_path) if model_7d_chart_path else None,
    "model_non7d_chart_path": str(model_non_7d_chart_path) if model_non_7d_chart_path else None,
}

summaries_7d = build_model_issue_table(df_7d, "7天无理由", detailed_dir_7d)
summaries_non7d = build_model_issue_table(df_non_7d, "非7天无理由", detailed_dir_non7d)

model_detail_map: Dict[str, Dict[str, Dict]] = {}

for entry in summaries_7d + summaries_non7d:
    model_name = entry.get("model")
    suffix = entry.get("suffix")
    if model_name not in model_detail_map:
        model_detail_map[model_name] = {}
    model_detail_map[model_name][suffix] = entry

summary_payload["model_details"] = model_detail_map


# 在主流程中调用生成分析报告的函数
generate_analysis_report(df, df_7d, df_non_7d, out_dir, start_date, end_date)

if args.generate_ppt:
    ppt_path = Path(args.ppt_path) if args.ppt_path else (out_dir / DEFAULT_PPT_PATH)
    llm_params = {
        "timeout": args.llm_timeout,
        "top_n": args.llm_top_n,
        "coverage": args.llm_coverage,
        "focus": args.llm_focus
    }
    try:
        generate_ppt(
            summary_payload,
            ppt_path,
            Path(args.ppt_template) if args.ppt_template else None,
            use_llm=args.use_llm,
            llm_params=llm_params
        )
    except (IOError, OSError) as exc:
        print(f"生成PPT失败（文件IO错误）: {exc}")
    except LLMGenerationError as exc:
        print(f"生成PPT失败（LLM调用错误）: {exc}")
    except Exception as exc:
        print(f"生成PPT失败（未知错误）: {exc}")
        import traceback
        traceback.print_exc()

print("✅ 所有处理完成，结果已保存到 output 目录！")
print("文件结构：")
print("output/")
print("├── 审核原因统计.xlsx")
print("├── 7天无理由_机型分布.xlsx")
print("├── 非7天无理由_机型分布.xlsx")
print("├── 审核原因占比.png")
print("├── 7天无理由_机型分布.png")
print("├── 非7天无理由_机型分布.png")
print("└── 详细数据/")
print("    ├── 7天无理由/")
print("    │   └── [机型名称]/")
print("    │       ├── [机型]_7天无理由_分类频次.xlsx")
print("    │       ├── [机型]_7天无理由_柱状图.png")
print("    │       └── [机型]_7天无理由_详细数据.xlsx")
print("    └── 非7天无理由/")
print("        └── [机型名称]/")
print("            ├── [机型]_非7天无理由_分类频次.xlsx")
print("            ├── [机型]_非7天无理由_柱状图.png")
print("            └── [机型]_非7天无理由_详细数据.xlsx")
if args.generate_ppt:
    print("├── report.pptx")
